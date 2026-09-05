import base64
import contextlib
import csv
import hashlib
import html as html_lib
import io
import logging
import os
import re
import shutil
import subprocess
import tempfile
import threading
import time
import gc
import json
from dataclasses import dataclass
from typing import Any, Generator
from concurrent.futures import ThreadPoolExecutor, as_completed

import fitz  # PyMuPDF
import requests
from docx import Document
from openpyxl import load_workbook
from openpyxl.utils import get_column_letter
from pptx import Presentation
from PIL import Image

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

@dataclass
class EngineConfig:
    vllm_base_url: str
    vllm_api_key: str
    text_model: str
    ocr_base_url: str = ""
    ocr_model: str = ""
    request_timeout: int = 600
    max_page_limit: int = 500
    max_concurrent_ocr: int = 96
    office_render_timeout: int = 600
    sliding_window_size: int = 50
    render_xlsx: bool = False
    render_docx: bool = True
    render_pptx: bool = True
    spreadsheet_chunk_rows: int = 200
    normalization_chunk_chars: int = 12000
    ocr_max_tokens: int = 8192

class DocprocEngine:
    def __init__(self, config: EngineConfig):
        self.config = config
        self._ocr_semaphore = threading.BoundedSemaphore(max(1, config.max_concurrent_ocr))

    def stream_extract(self, *, file_content: bytes, filename: str, page_limit: int | None = None, hint: str | None = None, prompt: str | None = None) -> Generator[str, None, None]:
        """Generator that sends heartbeat spaces to keep the network connection alive."""
        # Start the conversion in a separate thread so we can send heartbeats
        result_container = {}
        stop_heartbeat = threading.Event()

        def run_extraction():
            try:
                result_container['data'] = self.extract_document(
                    file_content=file_content, 
                    filename=filename, 
                    page_limit=page_limit,
                    hint=hint,
                    prompt=prompt,
                )
            except Exception as e:
                result_container['error'] = str(e)
            finally:
                stop_heartbeat.set()

        thread = threading.Thread(target=run_extraction)
        thread.start()

        while not stop_heartbeat.is_set():
            yield " " 
            time.sleep(5)

        thread.join()
        if 'error' in result_container:
            yield json.dumps({"error": result_container['error'], "status": "failed"})
        else:
            yield json.dumps(result_container.get('data', {}))

    def extract_document(self, *, file_content: bytes, filename: str, page_limit: int | None = None, hint: str | None = None, prompt: str | None = None) -> dict[str, Any]:
        result = self._extract_document_raw(
            file_content=file_content,
            filename=filename,
            page_limit=page_limit,
            hint=hint,
            prompt=prompt,
        )
        raw_text = str(result.get("raw_extracted_text") or result.get("normalized_text") or "").strip()
        if not raw_text or result.get("transcription_status") == "failed":
            return result
        try:
            normalized_text = self._normalize_extracted_text(raw_text, filename=filename)
        except Exception as exc:
            logger.warning("[%s] Text-model normalization failed: %s", filename, exc)
            result["normalized_text"] = raw_text
            result["quality_flags"] = [*(result.get("quality_flags") or []), "model_normalization_failed"]
            result["render_metadata"] = {
                **(result.get("render_metadata") or {}),
                "normalization_error": str(exc),
            }
            return result
        if normalized_text:
            result["normalized_text"] = normalized_text
            result["quality_flags"] = [*(result.get("quality_flags") or []), "text_model_normalized"]
        return result

    def _extract_document_raw(self, *, file_content: bytes, filename: str, page_limit: int | None = None, hint: str | None = None, prompt: str | None = None) -> dict[str, Any]:
        """Entry point for extraction."""
        ext = os.path.splitext(filename)[1].lower()
        
        # Priority: Request Limit > Centralized Default Limit
        limit = page_limit if page_limit else self.config.max_page_limit
        
        logger.info(f"--- START EXTRACTION: {filename} (Ext: {ext}) ---")
        try:
            if ext in {".png", ".jpg", ".jpeg", ".pdf"}:
                logger.info(f"[{filename}] Using multimodal text-model path")
                return self._extract_via_multimodal(file_content=file_content, filename=filename, page_limit=limit, hint=hint, prompt=prompt)
            
            # Office Speed Path
            text_export = ""
            should_render = True
            
            # 1. DOCX Path
            if ext in {".docx", ".doc", ".odt"}: 
                logger.info(f"[{filename}] Attempting Word/Doc text extraction...")
                text_export = self._extract_docx_text(file_content, limit)
                should_render = self.config.render_docx
                logger.info(f"[{filename}] Word extraction complete. Chars: {len(text_export)}, should_render: {should_render}")
                if ext == ".docx" and self._has_meaningful_text(text_export):
                    logger.info(f"[{filename}] SUCCESS: Returning native DOCX text.")
                    return self._build_result(
                        raw_text=text_export,
                        normalized_text=text_export,
                        quality_flags=["direct_text", "docx_native", "no_render"],
                    )

            # 2. PPTX Path
            elif ext in {".pptx", ".ppt", ".odp"}: 
                logger.info(f"[{filename}] Attempting PowerPoint/Slides text extraction...")
                text_export = self._extract_pptx_text(file_content, limit)
                should_render = self.config.render_pptx
                logger.info(f"[{filename}] PowerPoint extraction complete. Chars: {len(text_export)}, should_render: {should_render}")

            # 3. Excel/Tabular Path (STRENGTHENED)
            elif ext in {".xlsx", ".xls", ".xlsm", ".xlsb", ".csv", ".ods"}: 
                logger.info(f"[{filename}] Using complete spreadsheet extraction path.")
                spreadsheet = self._extract_spreadsheet_complete(file_content, filename=filename)
                if spreadsheet.get("normalized_text"):
                    return self._build_result(
                        raw_text=spreadsheet["raw_extracted_text"],
                        normalized_text=spreadsheet["normalized_text"],
                        quality_flags=spreadsheet["quality_flags"],
                        render_metadata=spreadsheet["render_metadata"],
                        structured_data=spreadsheet["structured_data"],
                    )
                if not self.config.render_xlsx:
                    return self._build_result(
                        raw_text="",
                        normalized_text="",
                        quality_flags=spreadsheet.get("quality_flags") or ["spreadsheet_structured_failed"],
                        render_metadata=spreadsheet.get("render_metadata") or {},
                        structured_data=spreadsheet.get("structured_data") or {},
                        error=spreadsheet.get("error") or "Spreadsheet extraction failed and whole-workbook rendering is disabled",
                        transcription_status="failed",
                    )
                text_export = ""
                should_render = True

            # 4. Outlook MSG Path
            elif ext == ".msg":
                logger.info(f"[{filename}] Using native Outlook MSG extraction path.")
                msg_result = self._extract_msg_complete(file_content, filename=filename)
                if msg_result.get("normalized_text"):
                    return self._build_result(
                        raw_text=msg_result["raw_extracted_text"],
                        normalized_text=msg_result["normalized_text"],
                        quality_flags=msg_result["quality_flags"],
                        render_metadata=msg_result["render_metadata"],
                        structured_data=msg_result["structured_data"],
                    )
                return self._build_result(
                    raw_text="",
                    normalized_text="",
                    quality_flags=msg_result.get("quality_flags") or ["msg_native_failed"],
                    render_metadata=msg_result.get("render_metadata") or {},
                    structured_data=msg_result.get("structured_data") or {},
                    error=msg_result.get("error") or "MSG extraction produced no readable content",
                    transcription_status="failed",
                )

            # Final Decision Logic
            if text_export.strip() and not should_render:
                logger.info(f"[{filename}] SUCCESS: Returning direct text (Rendering disabled or unnecessary).")
                return self._build_result(raw_text=text_export, normalized_text=text_export, quality_flags=["direct_text", "no_render"])

            if text_export.strip() and len(file_content) < 51200:
                logger.info(f"[{filename}] SUCCESS: Returning direct text (Small file optimization).")
                return self._build_result(raw_text=text_export, normalized_text=text_export, quality_flags=["direct_text", "small_file"])

            if not should_render and not text_export.strip():
                logger.warning(f"[{filename}] FAILED: Text extraction failed and rendering is disabled.")
                return self._build_result(raw_text="", normalized_text="", quality_flags=["failed_text_no_render"], error="Text extraction failed and rendering is disabled for this type")

            # 4. Fallback to rendering and the shared multimodal model
            logger.info(f"[{filename}] FALLBACK: Falling back to rendering path (should_render={should_render})")
            rendered = self._render_office_to_pdf_and_extract(file_content, filename, limit, hint=hint, prompt=prompt)
            if rendered:
                logger.info(f"[{filename}] Rendering success. Merging results.")
            else:
                logger.warning(f"[{filename}] Rendering failed (LibreOffice error).")
            
            return self._merge_extraction_results(rendered, text_export=text_export, route="render_plus_text", fallback_flag=f"{ext}_text_only")
        except Exception as e:
            logger.exception(f"Failure for {filename}")
            return self._build_result(raw_text="", normalized_text="", quality_flags=["crash"], error=str(e), transcription_status="failed")

    def _extract_via_multimodal(self, *, file_content: bytes, filename: str, page_limit: int | None, hint: str | None = None, prompt: str | None = None) -> dict[str, Any]:
        ext = os.path.splitext(filename)[1].lower()
        if ext in {".png", ".jpg", ".jpeg"}:
            img_b64 = self._optimize_and_encode(file_content)
            text = self._clean_model_text(self._multimodal_transcribe_page(img_b64, filename=filename, page_number=1, hint=hint, prompt=prompt))
            if not self._has_meaningful_text(text):
                return self._build_result(
                    raw_text="",
                    normalized_text="",
                    quality_flags=["multimodal_model", "empty_model_output"],
                    transcription_status="failed",
                    error="Multimodal extraction produced no readable content",
                )
            return self._build_result(raw_text=text, normalized_text=text, quality_flags=["multimodal_model"])

        pages_results = []
        try:
            with fitz.open(stream=file_content, filetype="pdf") as doc:
                limit = min(page_limit, len(doc)) if page_limit else len(doc)
                logger.info(f"[{filename}] Sliding window for {limit} pages...")
                
                # TUNABLE KNOB: window_size
                window_size = self.config.sliding_window_size
                for start_idx in range(0, limit, window_size):
                    end_idx = min(start_idx + window_size, limit)
                    batch_images = []
                    for i in range(start_idx, end_idx):
                        pix = doc[i].get_pixmap(matrix=fitz.Matrix(2.0, 2.0))
                        batch_images.append(self._optimize_and_encode(pix.tobytes("png")))
                        del pix
                    
                    batch_texts = [None] * len(batch_images)
                    with ThreadPoolExecutor(max_workers=self.config.max_concurrent_ocr) as executor:
                        future_to_idx = {
                            executor.submit(self._multimodal_transcribe_page, img, filename=filename, page_number=start_idx+i+1, hint=hint, prompt=prompt): i
                            for i, img in enumerate(batch_images)
                        }
                        for future in as_completed(future_to_idx):
                            idx = future_to_idx[future]
                            page_text = self._clean_model_text(future.result())
                            if self._has_meaningful_text(page_text):
                                batch_texts[idx] = f"--- {filename} (PAGE {start_idx+idx+1}) ---\n{page_text}"
                    
                    pages_results.extend(batch_texts)
                    batch_images.clear()
                    gc.collect()
                    logger.info(f"  Progress: {len(pages_results)}/{limit}")

        except Exception as e:
            logger.error(f"PDF failed: {e}")
            return self._build_result(raw_text="", normalized_text="", quality_flags=["failed"], error=str(e))

        full_text = "\n\n".join([p for p in pages_results if p]).strip()
        if not self._has_meaningful_text(full_text):
            return self._build_result(
                raw_text="",
                normalized_text="",
                quality_flags=["multimodal_model", "sliding_window", "empty_model_output"],
                transcription_status="failed",
                error="Multimodal extraction produced no readable content",
            )
        return self._build_result(raw_text=full_text, normalized_text=full_text, quality_flags=["multimodal_model", "sliding_window"])

    def _multimodal_transcribe_page(self, image_b64: str, *, filename: str, page_number: int, hint: str | None = None, prompt: str | None = None) -> str:
        headers = {"Content-Type": "application/json"}
        if self.config.vllm_api_key:
            headers["Authorization"] = f"Bearer {self.config.vllm_api_key}"
        
        base_url = (self.config.ocr_base_url or self.config.vllm_base_url).rstrip("/")
        if not base_url.endswith("/v1"): base_url = f"{base_url}/v1"

        base_prompt = prompt or "document parsing."
        requested_prompt = f"{hint}\n\n{base_prompt}" if hint else base_prompt
        ocr_model = self.config.ocr_model or self.config.text_model
        unlimited_ocr = "unlimited-ocr" in ocr_model.lower()
        # Unlimited-OCR has no chat template. Its literal <image> prefix and
        # per-request n-gram settings are required by the official recipe.
        final_prompt = requested_prompt
        if unlimited_ocr and not final_prompt.startswith("<image>"):
            final_prompt = f"<image>{final_prompt}"
        payload = {
            "model": ocr_model,
            "messages": [{"role": "user", "content": [
                {"type": "text", "text": final_prompt},
                {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{image_b64}"}}
            ]}],
            "temperature": 0.0,
            "max_tokens": self.config.ocr_max_tokens,
        }
        if unlimited_ocr:
            payload.update({
                "skip_special_tokens": False,
                "vllm_xargs": {"ngram_size": 35, "window_size": 128},
            })

        with self._ocr_semaphore:
            response = requests.post(
                f"{base_url}/chat/completions",
                headers=headers,
                json=payload,
                timeout=self.config.request_timeout,
            )
        response.raise_for_status()
        return response.json().get("choices", [{}])[0].get("message", {}).get("content", "").strip()

    def _normalize_extracted_text(self, text: str, *, filename: str) -> str:
        chunk_size = max(1000, self.config.normalization_chunk_chars)
        chunks = [text[start:start + chunk_size] for start in range(0, len(text), chunk_size)]
        normalized_chunks = []
        for index, chunk in enumerate(chunks, start=1):
            prompt = (
                "Normalize the extracted document text below into faithful Markdown. Preserve every fact, "
                "number, page marker, heading, and table value. Do not summarize, infer, or add commentary. "
                f"Document: {filename}. Part {index} of {len(chunks)}.\n\n{chunk}"
            )
            normalized = self._text_completion(prompt)
            normalized_chunks.append(normalized or chunk)
        return "\n\n".join(normalized_chunks).strip()

    def _text_completion(self, prompt: str) -> str:
        headers = {"Content-Type": "application/json"}
        if self.config.vllm_api_key:
            headers["Authorization"] = f"Bearer {self.config.vllm_api_key}"
        base_url = self.config.vllm_base_url.rstrip("/")
        if not base_url.endswith("/v1"):
            base_url = f"{base_url}/v1"
        response = requests.post(
            f"{base_url}/chat/completions",
            headers=headers,
            json={
                "model": self.config.text_model,
                "messages": [{"role": "user", "content": prompt}],
                "temperature": 0.0,
                "max_tokens": 4000,
                "chat_template_kwargs": {"enable_thinking": False},
            },
            timeout=self.config.request_timeout,
        )
        response.raise_for_status()
        return response.json().get("choices", [{}])[0].get("message", {}).get("content", "").strip()

    def _optimize_and_encode(self, img_bytes: bytes) -> str:
        img = Image.open(io.BytesIO(img_bytes))
        max_dim = 1600
        if max(img.size) > max_dim:
            ratio = max_dim / float(max(img.size))
            img = img.resize((int(img.size[0] * ratio), int(img.size[1] * ratio)), Image.LANCZOS)
        output = io.BytesIO()
        img.save(output, format="PNG", optimize=True)
        return base64.b64encode(output.getvalue()).decode("utf-8")

    def _render_office_to_pdf_and_extract(self, file_content: bytes, filename: str, page_limit: int | None, hint: str | None = None, prompt: str | None = None) -> dict[str, Any] | None:
        if not shutil.which("soffice"): return None
        ext = os.path.splitext(filename)[1].lower() or ".bin"
        with tempfile.TemporaryDirectory() as temp_dir:
            profile_dir = os.path.join(temp_dir, "profile")
            os.makedirs(profile_dir)
            in_p = os.path.join(temp_dir, "in" + ext)
            with open(in_p, "wb") as f: f.write(file_content)
            try:
                subprocess.run([
                    "soffice", f"-env:UserInstallation=file://{profile_dir}",
                    "--headless", "--convert-to", "pdf", "--outdir", temp_dir, in_p
                ], check=True, timeout=self.config.office_render_timeout)
                pdf_p = os.path.join(temp_dir, "in.pdf")
                if not os.path.exists(pdf_p):
                    pdf_p = os.path.join(temp_dir, os.path.splitext(os.path.basename(in_p))[0] + ".pdf")
                with open(pdf_p, "rb") as f: 
                    return self._extract_via_multimodal(file_content=f.read(), filename=filename, page_limit=page_limit, hint=hint, prompt=prompt)
            except: return None

    def _merge_extraction_results(self, rendered, text_export, route, fallback_flag):
        r_text = (rendered or {}).get("normalized_text", "").strip()
        t_text = (text_export or "").strip()
        if r_text and t_text:
            merged = f"{r_text}\n\n[STRUCTURED EXPORT]\n{t_text}"
            return self._build_result(raw_text=merged, normalized_text=merged, quality_flags=["merged"], render_metadata={"route": route})
        if r_text: return rendered
        if t_text: return self._build_result(raw_text=t_text, normalized_text=t_text, quality_flags=[fallback_flag])
        return self._build_result(raw_text="", normalized_text="", quality_flags=[fallback_flag, "failed"], transcription_status="failed")

    @staticmethod
    def _extract_docx_text(file_content: bytes, page_limit: int | None) -> str:
        try:
            doc = Document(io.BytesIO(file_content))
            parts = []
            for section in doc.sections:
                for paragraph in list(section.header.paragraphs) + list(section.footer.paragraphs):
                    text = paragraph.text.strip()
                    if text:
                        parts.append(text)
            for paragraph in doc.paragraphs:
                text = paragraph.text.strip()
                if text:
                    parts.append(text)
            for table_index, table in enumerate(doc.tables, start=1):
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                    if any(cells):
                        rows.append("\t".join(cells))
                if rows:
                    parts.append(f"--- DOCX TABLE {table_index} ---")
                    parts.extend(rows)
            return "\n".join(parts)
        except: return ""

    @staticmethod
    def _clean_model_text(text: str) -> str:
        cleaned = (text or "").strip()
        # Unlimited-OCR emits grounding labels and coordinate boxes. Preserve
        # the label text while removing coordinates from the Markdown output.
        cleaned = re.sub(r"<\|det\|>.*?<\|/det\|>", "", cleaned, flags=re.DOTALL)
        cleaned = re.sub(r"<\|ref\|>(.*?)<\|/ref\|>", r"\1", cleaned, flags=re.DOTALL)
        if cleaned.lower() in {"```markdown\n```", "```markdown```", "```", "``````"}:
            return ""
        if cleaned.startswith("```") and cleaned.endswith("```"):
            inner = cleaned.strip("`").strip()
            if inner.lower() in {"markdown", ""}:
                return ""
        return cleaned

    @staticmethod
    def _has_meaningful_text(text: str) -> bool:
        cleaned = DocprocEngine._clean_model_text(text)
        alnum_count = sum(ch.isalnum() for ch in cleaned)
        return alnum_count >= 3

    @staticmethod
    def _extract_pptx_text(file_content: bytes, page_limit: int | None) -> str:
        try:
            prs = Presentation(io.BytesIO(file_content))
            return "\n".join([s.shapes[i].text for s in prs.slides for i in range(len(s.shapes)) if hasattr(s.shapes[i], "text")])
        except: return ""

    def _extract_msg_complete(self, file_content: bytes, filename: str = "") -> dict[str, Any]:
        digest = hashlib.sha256(file_content).hexdigest()
        try:
            import extract_msg
        except Exception as exc:
            return {
                "raw_extracted_text": "",
                "normalized_text": "",
                "quality_flags": ["msg_native_failed", "missing_extract_msg_dependency"],
                "render_metadata": {
                    "route": "msg_native",
                    "content_sha256": digest,
                    "reader_error": str(exc),
                },
                "structured_data": {
                    "kind": "email_message",
                    "format": "msg",
                    "filename": filename,
                    "content_sha256": digest,
                },
                "error": f"extract-msg dependency unavailable: {exc}",
            }

        with tempfile.TemporaryDirectory() as temp_dir:
            msg_path = os.path.join(temp_dir, "message.msg")
            with open(msg_path, "wb") as f:
                f.write(file_content)

            msg = None
            try:
                msg = extract_msg.Message(msg_path)
                subject = self._clean_email_field(getattr(msg, "subject", None))
                sender = self._clean_email_field(getattr(msg, "sender", None) or getattr(msg, "senderEmail", None))
                to = self._clean_email_field(getattr(msg, "to", None))
                cc = self._clean_email_field(getattr(msg, "cc", None))
                bcc = self._clean_email_field(getattr(msg, "bcc", None))
                date = self._clean_email_field(getattr(msg, "date", None))
                message_id = self._clean_email_field(getattr(msg, "messageId", None))
                body = self._clean_email_field(getattr(msg, "body", None))
                html_body = getattr(msg, "htmlBody", None)
                html_text = self._html_to_text(html_body)
                body_text = body or html_text

                attachments = []
                for attachment in getattr(msg, "attachments", []) or []:
                    att_name = (
                        getattr(attachment, "longFilename", None)
                        or getattr(attachment, "shortFilename", None)
                        or getattr(attachment, "name", None)
                        or ""
                    )
                    att_size = None
                    data = getattr(attachment, "data", None)
                    if isinstance(data, (bytes, bytearray)):
                        att_size = len(data)
                    attachments.append({
                        "filename": self._clean_email_field(att_name),
                        "size_bytes": att_size,
                    })

                parts = [f"# OUTLOOK MSG: {filename}", f"SHA256: {digest}", ""]
                header_rows = [
                    ("Subject", subject),
                    ("From", sender),
                    ("To", to),
                    ("Cc", cc),
                    ("Bcc", bcc),
                    ("Date", date),
                    ("Message-ID", message_id),
                ]
                for label, value in header_rows:
                    if value:
                        parts.append(f"{label}: {value}")
                if attachments:
                    parts.append("")
                    parts.append("## Attachments")
                    for index, attachment in enumerate(attachments, start=1):
                        size = attachment.get("size_bytes")
                        size_text = f" ({size} bytes)" if size is not None else ""
                        parts.append(f"{index}. {attachment.get('filename') or 'Unnamed attachment'}{size_text}")
                if body_text:
                    parts.append("")
                    parts.append("## Body")
                    parts.append(body_text)

                normalized_text = "\n".join(parts).strip()
                meaningful_payload = "\n".join(value for value in (subject, sender, to, cc, bcc, date, body_text) if value)
                if not self._has_meaningful_text(meaningful_payload):
                    return {
                        "raw_extracted_text": "",
                        "normalized_text": "",
                        "quality_flags": ["msg_native", "empty_msg_output"],
                        "render_metadata": {
                            "route": "msg_native",
                            "content_sha256": digest,
                            "attachment_count": len(attachments),
                        },
                        "structured_data": {
                            "kind": "email_message",
                            "format": "msg",
                            "filename": filename,
                            "content_sha256": digest,
                            "subject": subject,
                            "sender": sender,
                            "to": to,
                            "cc": cc,
                            "bcc": bcc,
                            "date": date,
                            "message_id": message_id,
                            "attachments": attachments,
                        },
                        "error": "MSG extraction produced no readable subject/body/header text",
                    }

                return {
                    "raw_extracted_text": normalized_text,
                    "normalized_text": normalized_text,
                    "quality_flags": ["msg_native", "direct_text", "email_headers_extracted"],
                    "render_metadata": {
                        "route": "msg_native",
                        "content_sha256": digest,
                        "attachment_count": len(attachments),
                        "body_chars": len(body_text or ""),
                        "html_body_used": bool(html_text and not body),
                    },
                    "structured_data": {
                        "kind": "email_message",
                        "format": "msg",
                        "filename": filename,
                        "content_sha256": digest,
                        "subject": subject,
                        "sender": sender,
                        "to": to,
                        "cc": cc,
                        "bcc": bcc,
                        "date": date,
                        "message_id": message_id,
                        "attachments": attachments,
                    },
                }
            except Exception as exc:
                logger.warning(f"[{filename}] Native MSG extraction failed: {exc}")
                return {
                    "raw_extracted_text": "",
                    "normalized_text": "",
                    "quality_flags": ["msg_native_failed"],
                    "render_metadata": {
                        "route": "msg_native",
                        "content_sha256": digest,
                        "reader_error": str(exc),
                    },
                    "structured_data": {
                        "kind": "email_message",
                        "format": "msg",
                        "filename": filename,
                        "content_sha256": digest,
                    },
                    "error": str(exc),
                }
            finally:
                if msg is not None:
                    close = getattr(msg, "close", None)
                    if callable(close):
                        try:
                            close()
                        except Exception:
                            pass

    @staticmethod
    def _clean_email_field(value: Any) -> str:
        if value is None:
            return ""
        if isinstance(value, bytes):
            for encoding in ("utf-8", "utf-16", "latin-1"):
                try:
                    value = value.decode(encoding)
                    break
                except Exception:
                    continue
            if isinstance(value, bytes):
                value = value.decode("utf-8", errors="replace")
        return str(value).replace("\x00", "").strip()

    @classmethod
    def _html_to_text(cls, html_body: Any) -> str:
        html = cls._clean_email_field(html_body)
        if not html:
            return ""
        html = re.sub(r"(?is)<(script|style).*?>.*?</\1>", " ", html)
        html = re.sub(r"(?i)<br\s*/?>", "\n", html)
        html = re.sub(r"(?i)</p\s*>", "\n", html)
        html = re.sub(r"(?s)<[^>]+>", " ", html)
        html = html_lib.unescape(html)
        return "\n".join(line.strip() for line in html.splitlines() if line.strip())

    def _extract_spreadsheet_complete(self, file_content: bytes, filename: str = "") -> dict[str, Any]:
        ext = os.path.splitext(filename)[1].lower()
        try:
            if ext == ".csv":
                return self._extract_csv_complete(file_content, filename=filename)
            if ext in {".xlsx", ".xlsm"}:
                return self._extract_openpyxl_complete(file_content, filename=filename)
            return self._extract_pandas_spreadsheet_complete(file_content, filename=filename)
        except Exception as exc:
            logger.warning(f"[{filename}] Structured spreadsheet reader failed: {exc}")
            if ext in {".xls", ".ods"}:
                recovered = self._extract_legacy_spreadsheet_via_libreoffice(file_content, filename=filename, original_error=exc)
                if recovered:
                    return recovered
            digest = hashlib.sha256(file_content).hexdigest()
            return {
                "raw_extracted_text": "",
                "normalized_text": "",
                "quality_flags": ["spreadsheet_structured_failed"],
                "render_metadata": {
                    "route": "spreadsheet_structured",
                    "content_sha256": digest,
                    "reader_error": str(exc),
                    "rendered_sheet_count": 0,
                },
                "structured_data": {
                    "kind": "spreadsheet",
                    "filename": filename,
                    "content_sha256": digest,
                    "sheets": [],
                    "chunks": [],
                },
                "error": str(exc),
            }

    def _extract_csv_complete(self, file_content: bytes, filename: str = "") -> dict[str, Any]:
        decoded = file_content.decode("utf-8-sig", errors="replace")
        rows = list(csv.reader(io.StringIO(decoded)))
        chunk_rows = max(1, self.config.spreadsheet_chunk_rows)
        chunks = []
        text_parts = [f"# CSV: {filename}", ""]
        for start in range(0, len(rows), chunk_rows):
            end = min(start + chunk_rows, len(rows))
            lines = [",".join(row) for row in rows[start:end]]
            chunk_text = "\n".join(lines)
            chunks.append({
                "text": chunk_text,
                "metadata": {
                    "chunk_kind": "spreadsheet_range",
                    "sheet_name": "CSV",
                    "row_start": start + 1,
                    "row_end": end,
                    "column_start": "A",
                    "column_end": get_column_letter(max((len(row) for row in rows[start:end]), default=1)),
                },
            })
            text_parts.append(f"## SHEET: CSV rows {start + 1}-{end}")
            text_parts.append(chunk_text)
            text_parts.append("")
        digest = hashlib.sha256(file_content).hexdigest()
        structured_data = {
            "kind": "spreadsheet",
            "format": "csv",
            "filename": filename,
            "content_sha256": digest,
            "sheets": [{"name": "CSV", "row_count": len(rows), "column_count": max((len(row) for row in rows), default=0)}],
            "chunks": chunks,
        }
        return {
            "raw_extracted_text": "\n".join(text_parts).strip(),
            "normalized_text": "\n".join(text_parts).strip(),
            "quality_flags": ["spreadsheet_structured", "raw_complete", "artifact_backed", "csv"],
            "render_metadata": {
                "route": "spreadsheet_structured",
                "content_sha256": digest,
                "sheet_count": 1,
                "row_count": len(rows),
                "rendered_sheet_count": 0,
            },
            "structured_data": structured_data,
        }

    def _extract_openpyxl_complete(self, file_content: bytes, filename: str = "") -> dict[str, Any]:
        ext = os.path.splitext(filename)[1].lower()
        digest = hashlib.sha256(file_content).hexdigest()
        formula_wb = load_workbook(io.BytesIO(file_content), data_only=False, read_only=True)
        value_wb = load_workbook(io.BytesIO(file_content), data_only=True, read_only=True)
        chunk_rows = max(1, self.config.spreadsheet_chunk_rows)
        text_parts = [f"# WORKBOOK: {filename}", f"SHA256: {digest}", ""]
        sheets = []
        chunks = []
        flags = ["spreadsheet_structured", "raw_complete", "artifact_backed"]

        for sheet_name in formula_wb.sheetnames:
            ws_formula = formula_wb[sheet_name]
            ws_value = value_wb[sheet_name]
            max_row = ws_formula.max_row or 0
            max_col = ws_formula.max_column or 0
            column_end = get_column_letter(max_col) if max_col else "A"
            hidden = ws_formula.sheet_state != "visible"
            merged_cells = getattr(ws_formula, "merged_cells", None)
            merged_ranges = [str(rng) for rng in getattr(merged_cells, "ranges", [])]
            sheet_info = {
                "name": sheet_name,
                "row_count": max_row,
                "column_count": max_col,
                "column_start": "A",
                "column_end": column_end,
                "hidden": hidden,
                "merged_ranges": merged_ranges,
            }
            sheets.append(sheet_info)
            if hidden and "hidden_sheets_present" not in flags:
                flags.append("hidden_sheets_present")
            if merged_ranges and "merged_cells_present" not in flags:
                flags.append("merged_cells_present")

            text_parts.append(f"## SHEET: {sheet_name}")
            text_parts.append(f"ROWS: {max_row} COLUMNS: {max_col} HIDDEN: {hidden}")
            if merged_ranges:
                text_parts.append("MERGED_RANGES: " + ", ".join(merged_ranges))

            chunk_buffer = []
            chunk_start = 1
            for row_index, (formula_row, value_row) in enumerate(
                zip(ws_formula.iter_rows(values_only=False), ws_value.iter_rows(values_only=True)),
                start=1,
            ):
                values = []
                for col_index in range(1, max_col + 1):
                    formula_cell = formula_row[col_index - 1] if col_index <= len(formula_row) else None
                    cached_value = value_row[col_index - 1] if col_index <= len(value_row) else None
                    raw_value = formula_cell.value if formula_cell is not None else None
                    formatted = self._format_cell_value(raw_value, cached_value)
                    values.append(formatted)
                    if isinstance(raw_value, str) and raw_value.startswith("=") and "formulas_present" not in flags:
                        flags.append("formulas_present")
                row_text = f"{row_index}\t" + "\t".join(values)
                text_parts.append(row_text)
                chunk_buffer.append(row_text)

                if len(chunk_buffer) >= chunk_rows:
                    chunks.append({
                        "text": "\n".join(chunk_buffer),
                        "metadata": {
                            "chunk_kind": "spreadsheet_range",
                            "sheet_name": sheet_name,
                            "row_start": chunk_start,
                            "row_end": row_index,
                            "column_start": "A",
                            "column_end": column_end,
                        },
                    })
                    chunk_buffer = []
                    chunk_start = row_index + 1

            if chunk_buffer:
                chunks.append({
                    "text": "\n".join(chunk_buffer),
                    "metadata": {
                        "chunk_kind": "spreadsheet_range",
                        "sheet_name": sheet_name,
                        "row_start": chunk_start,
                        "row_end": max_row,
                        "column_start": "A",
                        "column_end": column_end,
                    },
                })
            text_parts.append("")

        formula_wb.close()
        value_wb.close()
        structured_data = {
            "kind": "spreadsheet",
            "format": ext.lstrip("."),
            "filename": filename,
            "content_sha256": digest,
            "sheets": sheets,
            "chunks": chunks,
        }
        normalized_text = "\n".join(text_parts).strip()
        return {
            "raw_extracted_text": normalized_text,
            "normalized_text": normalized_text,
            "quality_flags": flags,
            "render_metadata": {
                "route": "spreadsheet_structured",
                "content_sha256": digest,
                "sheet_count": len(sheets),
                "sheets": sheets,
                "spreadsheet_chunk_count": len(chunks),
                "rendered_sheet_count": 0,
            },
            "structured_data": structured_data,
        }

    def _extract_pandas_spreadsheet_complete(self, file_content: bytes, filename: str = "") -> dict[str, Any]:
        import pandas as pd
        ext = os.path.splitext(filename)[1].lower()
        engine = "pyxlsb" if ext == ".xlsb" else "xlrd" if ext == ".xls" else None
        digest = hashlib.sha256(file_content).hexdigest()
        with contextlib.redirect_stdout(io.StringIO()), contextlib.redirect_stderr(io.StringIO()):
            sheets_map = pd.read_excel(io.BytesIO(file_content), sheet_name=None, engine=engine, header=None, dtype=object)
        chunk_rows = max(1, self.config.spreadsheet_chunk_rows)
        text_parts = [f"# WORKBOOK: {filename}", f"SHA256: {digest}", ""]
        sheets = []
        chunks = []
        for sheet_name, df in sheets_map.items():
            row_count = int(df.shape[0])
            column_count = int(df.shape[1])
            column_end = get_column_letter(column_count) if column_count else "A"
            sheets.append({
                "name": sheet_name,
                "row_count": row_count,
                "column_count": column_count,
                "column_start": "A",
                "column_end": column_end,
                "hidden": None,
                "merged_ranges": [],
            })
            text_parts.append(f"## SHEET: {sheet_name}")
            text_parts.append(f"ROWS: {row_count} COLUMNS: {column_count}")
            rows = []
            for idx, row in df.iterrows():
                values = ["" if pd.isna(value) else str(value) for value in row.tolist()]
                rows.append(f"{idx + 1}\t" + "\t".join(values))
            for start in range(0, len(rows), chunk_rows):
                end = min(start + chunk_rows, len(rows))
                chunk_text = "\n".join(rows[start:end])
                chunks.append({
                    "text": chunk_text,
                    "metadata": {
                        "chunk_kind": "spreadsheet_range",
                        "sheet_name": sheet_name,
                        "row_start": start + 1,
                        "row_end": end,
                        "column_start": "A",
                        "column_end": column_end,
                    },
                })
                text_parts.append(chunk_text)
            text_parts.append("")
        structured_data = {
            "kind": "spreadsheet",
            "format": ext.lstrip("."),
            "filename": filename,
            "content_sha256": digest,
            "sheets": sheets,
            "chunks": chunks,
        }
        normalized_text = "\n".join(text_parts).strip()
        return {
            "raw_extracted_text": normalized_text,
            "normalized_text": normalized_text,
            "quality_flags": ["spreadsheet_structured", "raw_complete", "artifact_backed"],
            "render_metadata": {
                "route": "spreadsheet_structured",
                "content_sha256": digest,
                "sheet_count": len(sheets),
                "sheets": sheets,
                "spreadsheet_chunk_count": len(chunks),
                "rendered_sheet_count": 0,
            },
            "structured_data": structured_data,
        }

    def _extract_legacy_spreadsheet_via_libreoffice(self, file_content: bytes, filename: str, original_error: Exception) -> dict[str, Any] | None:
        if not shutil.which("soffice"):
            logger.warning(f"[{filename}] LibreOffice is unavailable for legacy spreadsheet recovery.")
            return None
        ext = os.path.splitext(filename)[1].lower() or ".xls"
        with tempfile.TemporaryDirectory() as temp_dir:
            profile_dir = os.path.join(temp_dir, "profile")
            os.makedirs(profile_dir)
            in_path = os.path.join(temp_dir, "in" + ext)
            with open(in_path, "wb") as f:
                f.write(file_content)
            try:
                subprocess.run(
                    [
                        "soffice",
                        f"-env:UserInstallation=file://{profile_dir}",
                        "--headless",
                        "--convert-to",
                        "xlsx",
                        "--outdir",
                        temp_dir,
                        in_path,
                    ],
                    check=True,
                    timeout=self.config.office_render_timeout,
                    stdout=subprocess.DEVNULL,
                    stderr=subprocess.DEVNULL,
                )
                converted_path = os.path.join(temp_dir, "in.xlsx")
                if not os.path.exists(converted_path):
                    matches = [p for p in os.listdir(temp_dir) if p.lower().endswith(".xlsx")]
                    if not matches:
                        return None
                    converted_path = os.path.join(temp_dir, matches[0])
                with open(converted_path, "rb") as f:
                    recovered = self._extract_openpyxl_complete(f.read(), filename=f"{filename}.converted.xlsx")
                recovered["quality_flags"] = list(dict.fromkeys((recovered.get("quality_flags") or []) + ["legacy_spreadsheet_recovered"]))
                recovered.setdefault("render_metadata", {})
                recovered["render_metadata"]["route"] = "legacy_spreadsheet_converted_to_xlsx"
                recovered["render_metadata"]["legacy_reader_error"] = str(original_error)
                recovered["render_metadata"]["original_filename"] = filename
                recovered.setdefault("structured_data", {})
                recovered["structured_data"]["original_filename"] = filename
                recovered["structured_data"]["converted_format"] = "xlsx"
                return recovered
            except Exception as exc:
                logger.warning(f"[{filename}] LibreOffice legacy spreadsheet recovery failed: {exc}")
                return None

    @staticmethod
    def _format_cell_value(raw_value: Any, cached_value: Any) -> str:
        if raw_value is None and cached_value is None:
            return ""
        if isinstance(raw_value, str) and raw_value.startswith("="):
            if cached_value is not None and cached_value != raw_value:
                return f"{raw_value} => {cached_value}"
            return raw_value
        return str(raw_value if raw_value is not None else cached_value)

    @staticmethod
    def _extract_xlsx_text(file_content: bytes, page_limit: int | None, filename: str = "") -> tuple[str, list[str]]:
        """
        High-fidelity Excel/CSV extraction.
        Converts sheets to Markdown tables for optimal LLM consumption.
        """
        ext = os.path.splitext(filename)[1].lower() if filename else ""
        try:
            import pandas as pd
            excel_file = io.BytesIO(file_content)
            
            # Special Case: CSV
            if ext == ".csv":
                logger.info(f"Attempting pandas CSV extraction...")
                df = pd.read_csv(excel_file)
                if df.empty: return "", []
                df = df.dropna(how='all').dropna(axis=1, how='all')
                text_out = df.to_markdown(index=False)
                return text_out, ["CSV_Sheet"]

            # Standard Excel Path
            logger.info(f"Attempting pandas extraction for Excel ({ext})...")
            # engine='openpyxl' supports xlsx, xlsm, xltx, xltm. 
            # engine='pyxlsb' for xlsb. 
            # engine='xlrd' for old xls.
            
            engine = 'openpyxl'
            if ext == ".xlsb": engine = 'pyxlsb'
            elif ext == ".xls": engine = 'xlrd'

            all_sheets = pd.read_excel(excel_file, sheet_name=None, engine=engine)
            
            res = []
            sheet_names = []
            
            for name, df in all_sheets.items():
                sheet_names.append(name)
                if df.empty:
                    continue
                
                # Clean up: remove entirely empty rows/columns to save tokens
                df = df.dropna(how='all').dropna(axis=1, how='all')
                if df.empty:
                    continue

                res.append(f"### SHEET: {name}")
                # Convert to high-contrast Markdown table
                try:
                    res.append(df.to_markdown(index=False))
                except Exception as table_err:
                    logger.warning(f"Markdown table conversion failed for sheet {name}, using TSV: {table_err}")
                    res.append(df.to_csv(sep='\t', index=False))
                res.append("\n")

            text_out = "\n\n".join(res)
            logger.info(f"Pandas extraction complete. Extracted {len(text_out)} characters from {len(sheet_names)} sheets.")
            return text_out, sheet_names
        except Exception as e:
            logger.error(f"Advanced pandas Excel extraction failed: {e}")
            try:
                # Basic fallback if pandas fails (openpyxl only supports modern XML formats)
                logger.info("Attempting openpyxl read_only fallback...")
                from openpyxl import load_workbook
                wb = load_workbook(io.BytesIO(file_content), data_only=True, read_only=True)
                res = []
                for name in wb.sheetnames:
                    rows = []
                    # Limit rows in fallback to prevent massive strings
                    for row_idx, r in enumerate(wb[name].iter_rows(values_only=True)):
                        rows.append("\t".join([str(c) if c else "" for c in r]))
                        if row_idx > 5000: # Safety cap
                            rows.append("... [TRUNCATED DUE TO SIZE] ...")
                            break
                    res.append(f"--- {name} ---\n" + "\n".join(rows))
                
                text_out = "\n".join(res)
                logger.info(f"Openpyxl fallback complete. Extracted {len(text_out)} characters.")
                return text_out, wb.sheetnames
            except Exception as e2:
                logger.error(f"Openpyxl fallback also failed: {e2}")
                return "", []

    @staticmethod
    def _build_result(*, raw_text: str, normalized_text: str, quality_flags: list[str], render_metadata: dict = None, structured_data: dict = None, transcription_status="complete", error=None) -> dict:
        return {"raw_extracted_text": raw_text, "normalized_text": normalized_text, "extraction_mode": "docproc_remote", "transcription_status": transcription_status, "quality_flags": quality_flags, "render_metadata": render_metadata or {}, "structured_data": structured_data or {}, "error": error}
